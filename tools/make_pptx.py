"""Generate the Babel's Glyph pitch deck as a .pptx file.

Theme is drawn directly from the in-game palette (game/render.py):

  * Sandstone / dawn        zone 1 — warm gold + peach
  * Forge crimson + ember   zone 2 — deep red, hot orange
  * Sky workshop            zone 3 — cool blue + cream
  * Parchment + ink         narration / tech slides
  * Glyph-gold              global accent

Every slide has a hand-painted parchment-grain background, illuminated-
manuscript corner ornaments, and a glyph-circle sigil that matches the
collectible drawn in-game (world.py:_glyph()).

Run:  python tools/make_pptx.py
Output: BabelsGlyph_Presentation.pptx (10 slides, 16:9).
"""
from __future__ import annotations

import io
import math
import random
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Palette — pulled from game/render.py to match the live game
# ---------------------------------------------------------------------------
# Sandstone zone (dawn)
SANDSTONE     = (216, 179, 117)
SAND_LIGHT    = (236, 206, 152)
SAND_DARK     = (168, 132, 80)
SKY_DAWN_TOP  = (252, 218, 158)
SKY_DAWN_BOT  = (210, 130, 90)

# Forge zone (crimson + ember)
SKY_FORGE_TOP = ( 90,  50,  60)
SKY_FORGE_BOT = (200,  80,  50)
EMBER         = (255, 122,  43)
EMBER_DIM     = (180,  70,  20)

# Sky workshop zone (cool dawn cream)
SKY_WORK_TOP  = (170, 200, 220)
SKY_WORK_BOT  = (220, 220, 200)

# Universal accents
GLYPH_GOLD    = (255, 210, 110)
GLYPH_GLOW_S  = (255, 244, 180)
LAPIS         = ( 28,  74, 138)
LAPIS_LIGHT   = ( 60, 120, 200)
COPPER        = (184, 115,  51)
COPPER_LIGHT  = (220, 160,  90)
PARCHMENT     = (228, 208, 160)
PARCHMENT_HI  = (244, 230, 196)
INK           = ( 38,  26,  18)         # near-black brown
INK_SOFT      = ( 70,  52,  38)
INK_FAINT     = (110,  88,  68)

FONT_TITLE = "Cinzel"          # roman-capital serif (falls back gracefully)
FONT_TITLE_FB = "Trajan Pro"
FONT_BODY  = "Segoe UI"
FONT_MONO  = "Consolas"

OUT_PATH = Path(__file__).resolve().parent.parent / "BabelsGlyph_Presentation.pptx"
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
PX_PER_IN  = 144     # background image resolution


def rgb(c):
    return RGBColor(*c)


# ---------------------------------------------------------------------------
# Background painter — produces a parchment / sky / forge gradient
#   with subtle grain, vignette, and corner ornaments.
# ---------------------------------------------------------------------------
def _vlerp(a, b, t):
    return tuple(int(a[i] + (b[i] - a[i]) * t) for i in range(3))


def _add_grain(img: Image.Image, strength: int = 8, seed: int = 7) -> Image.Image:
    """Sprinkle warm noise on top of the gradient for parchment feel."""
    w, h = img.size
    rnd = random.Random(seed)
    noise = Image.new("RGB", (w, h), (0, 0, 0))
    px = noise.load()
    for y in range(h):
        for x in range(w):
            n = rnd.randint(-strength, strength)
            px[x, y] = (max(0, n), max(0, n), max(0, n // 2))
    noise = noise.filter(ImageFilter.GaussianBlur(0.6))
    return Image.blend(img, Image.eval(noise, lambda v: v + 0).convert("RGB"),
                       0.10)


def _vignette(img: Image.Image, depth: float = 0.45) -> Image.Image:
    w, h = img.size
    mask = Image.new("L", (w, h), 0)
    d = ImageDraw.Draw(mask)
    cx, cy = w / 2, h / 2
    max_r = math.hypot(cx, cy)
    steps = 80
    for i in range(steps, 0, -1):
        r = max_r * (i / steps)
        a = int(255 * (1 - i / steps) * depth)
        d.ellipse((cx - r, cy - r, cx + r, cy + r), fill=a)
    dark = Image.new("RGB", (w, h), (0, 0, 0))
    return Image.composite(dark, img, mask)


def _draw_glyph_sigil(d: ImageDraw.ImageDraw, cx, cy, r,
                      color, alpha=180, line_w=2):
    """The collectible glyph drawn in-game: circle + cross + diagonal."""
    col = (*color, alpha) if len(color) == 3 else color
    # Outer ring
    d.ellipse((cx - r, cy - r, cx + r, cy + r),
              outline=col, width=line_w)
    # Inner ring
    d.ellipse((cx - r * 0.55, cy - r * 0.55,
               cx + r * 0.55, cy + r * 0.55),
              outline=col, width=max(1, line_w - 1))
    # Cross + diagonal arms
    d.line((cx - r, cy, cx + r, cy), fill=col, width=line_w)
    d.line((cx, cy - r, cx, cy + r), fill=col, width=line_w)
    d.line((cx - r * 0.7, cy - r * 0.7,
            cx + r * 0.7, cy + r * 0.7), fill=col, width=line_w)
    d.line((cx - r * 0.7, cy + r * 0.7,
            cx + r * 0.7, cy - r * 0.7), fill=col, width=line_w)


def _draw_corner_ornament(d: ImageDraw.ImageDraw, x, y, size, color,
                          flip_x=False, flip_y=False, w=2):
    """Illuminated-manuscript corner flourish: bracket + curl + dot."""
    # Anchor coords — treat (x,y) as the inner tip of the corner.
    sx = -1 if flip_x else 1
    sy = -1 if flip_y else 1
    # Long bracket
    d.line((x, y, x + sx * size, y), fill=color, width=w)
    d.line((x, y, x, y + sy * size), fill=color, width=w)
    # Inner shorter bracket
    pad = size * 0.15
    d.line((x + sx * pad, y + sy * pad,
            x + sx * size * 0.7, y + sy * pad), fill=color, width=w)
    d.line((x + sx * pad, y + sy * pad,
            x + sx * pad, y + sy * size * 0.7), fill=color, width=w)
    # Decorative dot at the inner corner
    r = max(2, w + 1)
    d.ellipse((x + sx * pad - r, y + sy * pad - r,
               x + sx * pad + r, y + sy * pad + r), fill=color)
    # Outer terminal dot
    r2 = max(1, w)
    d.ellipse((x + sx * size - r2, y - r2,
               x + sx * size + r2, y + r2), fill=color)
    d.ellipse((x - r2, y + sy * size - r2,
               x + r2, y + sy * size + r2), fill=color)


def _gradient_bg(top, bot, w, h):
    img = Image.new("RGB", (w, h), top)
    px = img.load()
    for y in range(h):
        t = y / max(1, h - 1)
        col = _vlerp(top, bot, t)
        for x in range(w):
            px[x, y] = col
    return img


def make_bg_image(theme: str) -> bytes:
    """Render a full-bleed background PNG sized to the slide.

    Themes:
      title    -> dawn sandstone gradient with bright glyph
      ink      -> deep ink-brown w/ embers (used for tech / pitch / numbers)
      forge    -> dark crimson with ember pulse
      sky      -> cool dawn cream (used for roadmap / world)
      parchment -> warm cream parchment (used for pillars / movement)
    """
    w = int(SLIDE_W_IN * PX_PER_IN)
    h = int(SLIDE_H_IN * PX_PER_IN)

    if theme == "title":
        img = _gradient_bg(SKY_DAWN_TOP, SKY_DAWN_BOT, w, h)
    elif theme == "forge":
        img = _gradient_bg(SKY_FORGE_TOP, SKY_FORGE_BOT, w, h)
    elif theme == "sky":
        img = _gradient_bg(SKY_WORK_TOP, SKY_WORK_BOT, w, h)
    elif theme == "parchment":
        img = _gradient_bg(PARCHMENT_HI, PARCHMENT, w, h)
    else:  # "ink"
        img = _gradient_bg((28, 20, 14), (14, 10, 8), w, h)

    # Glyph sigil watermark — large, faint, off-center on the right.
    over = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    od = ImageDraw.Draw(over)
    if theme == "title":
        _draw_glyph_sigil(od, int(w * 0.78), int(h * 0.46),
                          int(h * 0.32), GLYPH_GOLD, alpha=110, line_w=4)
    elif theme == "ink":
        _draw_glyph_sigil(od, int(w * 0.86), int(h * 0.18),
                          int(h * 0.13), GLYPH_GOLD, alpha=55, line_w=2)
        _draw_glyph_sigil(od, int(w * 0.10), int(h * 0.86),
                          int(h * 0.10), EMBER,      alpha=40, line_w=2)
    elif theme == "forge":
        _draw_glyph_sigil(od, int(w * 0.84), int(h * 0.28),
                          int(h * 0.18), GLYPH_GOLD, alpha=85, line_w=3)
    elif theme == "sky":
        _draw_glyph_sigil(od, int(w * 0.82), int(h * 0.78),
                          int(h * 0.16), LAPIS,      alpha=55, line_w=2)
    else:  # parchment
        _draw_glyph_sigil(od, int(w * 0.86), int(h * 0.85),
                          int(h * 0.13), INK_SOFT,   alpha=60, line_w=2)

    img = Image.alpha_composite(img.convert("RGBA"), over).convert("RGB")

    # Subtle parchment grain everywhere except true ink (kept clean for text)
    if theme != "ink":
        img = _add_grain(img, strength=6, seed=hash(theme) & 0xFFF)
    else:
        img = _add_grain(img, strength=4, seed=2)

    # Vignette pulls focus to the centre
    img = _vignette(img, depth=0.55 if theme == "ink" else 0.30)

    # Corner ornaments — small brass brackets in each corner
    final = Image.alpha_composite(
        img.convert("RGBA"), Image.new("RGBA", (w, h), (0, 0, 0, 0)))
    fd = ImageDraw.Draw(final)
    pad = int(h * 0.05)
    orn_size = int(h * 0.06)
    if theme in ("title", "parchment", "sky"):
        orn_color = (*INK_SOFT, 220)
    elif theme == "forge":
        orn_color = (*GLYPH_GOLD, 230)
    else:
        orn_color = (*GLYPH_GOLD, 210)
    _draw_corner_ornament(fd, pad, pad, orn_size, orn_color, w=3)
    _draw_corner_ornament(fd, w - pad, pad, orn_size, orn_color,
                          flip_x=True, w=3)
    _draw_corner_ornament(fd, pad, h - pad, orn_size, orn_color,
                          flip_y=True, w=3)
    _draw_corner_ornament(fd, w - pad, h - pad, orn_size, orn_color,
                          flip_x=True, flip_y=True, w=3)

    final = final.convert("RGB")
    buf = io.BytesIO()
    final.save(buf, format="PNG", optimize=True)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------
def add_bg(slide, theme: str):
    sw = slide.part.package.presentation_part.presentation.slide_width
    sh = slide.part.package.presentation_part.presentation.slide_height
    png = make_bg_image(theme)
    pic = slide.shapes.add_picture(io.BytesIO(png), 0, 0,
                                   width=sw, height=sh)
    # Send to back: pptx puts shapes on top in z-order they're added,
    # so as long as we add bg first, it stays at the bottom.
    return pic


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=PARCHMENT, font=FONT_BODY, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, italic=False, spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing is not None:
        p.line_spacing = spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return tb


def add_paragraphs(slide, x, y, w, h, lines, *, size=16, color=PARCHMENT,
                   font=FONT_BODY, bullet=True, line_spacing=1.25,
                   bullet_glyph="◆"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, c, bold = item
        else:
            text, c, bold = item, color, False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        if bullet:
            br = p.add_run()
            br.text = f"{bullet_glyph}  "
            br.font.name = font
            br.font.size = Pt(size)
            br.font.color.rgb = rgb(GLYPH_GOLD)
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(c)
    return tb


def add_rule(slide, x, y, w, color=GLYPH_GOLD, thickness=Pt(2)):
    line = slide.shapes.add_connector(1, x, y, x + w, y)
    line.line.color.rgb = rgb(color)
    line.line.width = thickness
    return line


def add_double_rule(slide, x, y, w, color=GLYPH_GOLD):
    """Illuminated-manuscript double rule: thick + thin separated."""
    add_rule(slide, x, y, w, color=color, thickness=Pt(2.5))
    add_rule(slide, x, y + Inches(0.10), Inches(w / 914400 * 0.45),
             color=color, thickness=Pt(0.75))


def add_panel(slide, x, y, w, h, *, fill=None, fill_alpha=None,
              border=GLYPH_GOLD, border_w=Pt(1.25)):
    """Parchment-style panel — translucent fill + brass-thin border."""
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        panel.fill.background()
    else:
        panel.fill.solid()
        panel.fill.fore_color.rgb = rgb(fill)
    panel.line.color.rgb = rgb(border)
    panel.line.width = border_w
    panel.shadow.inherit = False
    # Decorative inner stroke
    pad = Inches(0.06)
    inner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x + pad, y + pad,
                                   w - pad * 2, h - pad * 2)
    inner.fill.background()
    inner.line.color.rgb = rgb(border)
    inner.line.width = Pt(0.4)
    inner.shadow.inherit = False
    return panel


def add_kicker(slide, x, y, label, color=GLYPH_GOLD):
    """Small uppercase eyebrow label with bullet flourishes."""
    add_text(slide, x, y, Inches(8), Inches(0.4),
             f"◆  {label}  ◆",
             size=11, color=color, bold=True, font=FONT_BODY)


def add_footer(slide, page, total, color=INK_FAINT):
    sw = slide.part.package.presentation_part.presentation.slide_width
    sh = slide.part.package.presentation_part.presentation.slide_height
    add_text(slide, Inches(0.6), sh - Inches(0.42), Inches(6), Inches(0.3),
             "BABEL'S GLYPH  ·  PITCH DECK",
             size=9, color=color, bold=True)
    add_text(slide, sw - Inches(1.5), sh - Inches(0.42), Inches(0.9),
             Inches(0.3), f"{page:02d} / {total:02d}",
             size=9, color=color, bold=True, align=PP_ALIGN.RIGHT,
             font=FONT_MONO)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]
    SW = prs.slide_width
    SH = prs.slide_height
    TOTAL = 11

    # ============================================================ 1. TITLE
    s = prs.slides.add_slide(blank)
    add_bg(s, "title")

    add_text(s, Inches(0.9), Inches(0.95), Inches(8), Inches(0.5),
             "AN ENDLESS RUNNER  ·  BUILT IN PYTHON",
             size=14, color=INK_SOFT, bold=True, font=FONT_BODY)

    # Title set in serif caps to match the in-game UI.
    add_text(s, Inches(0.9), Inches(1.45), Inches(11.5), Inches(2.6),
             "BABEL'S GLYPH",
             size=120, bold=True, color=INK, font=FONT_TITLE)

    add_double_rule(s, Inches(0.95), Inches(4.05), Inches(2.6),
                    color=COPPER)

    add_text(s, Inches(0.9), Inches(4.30), Inches(11.5), Inches(0.8),
             "Outrun history itself.",
             size=34, color=INK_SOFT, font=FONT_TITLE, italic=True)

    add_text(s, Inches(0.9), Inches(5.20), Inches(11), Inches(1.4),
             "A hand-crafted 2-D platformer about a glyph-thief sprinting "
             "through three\n"
             "collapsing zones — sandstone ruins, da Vinci's forge, and "
             "the workshop above the clouds.",
             size=18, color=INK_SOFT)

    add_text(s, Inches(0.9), SH - Inches(0.85), Inches(11.5), Inches(0.4),
             "Game  ·  Features  ·  Tech  ·  Future Scope",
             size=12, color=INK_FAINT, bold=True, font=FONT_BODY)

    # ============================================================ 2. PITCH
    s = prs.slides.add_slide(blank)
    add_bg(s, "ink")
    add_kicker(s, Inches(0.7), Inches(0.55), "THE PITCH")
    add_text(s, Inches(0.7), Inches(0.95), Inches(11), Inches(1.0),
             "What Is Babel's Glyph?",
             size=42, bold=True, color=PARCHMENT_HI, font=FONT_TITLE)
    add_double_rule(s, Inches(0.75), Inches(1.95), Inches(2.2))

    add_text(s, Inches(0.7), Inches(2.25), Inches(12), Inches(2.0),
             "An endless side-scrolling runner where the world literally "
             "falls apart\n"
             "behind you. Sprint, dash, wall-run, and hurl glyph-bombs "
             "through three\n"
             "hand-stitched zones — every run procedurally remixed, "
             "never the same twice.",
             size=22, color=PARCHMENT, spacing=1.3)

    panels = [
        ("FAST",  "Locked 60 FPS. Coyote time, jump-buffer, dash i-frames.",
         GLYPH_GOLD),
        ("DEEP",  "Double-jump · dash · wall-run · slide · glyph-bomb.",
         LAPIS_LIGHT),
        ("ALIVE", "Procedural world from authored chunks. Emergent difficulty.",
         EMBER),
    ]
    px = Inches(0.7)
    py = Inches(4.65)
    pw = Inches(4.0)
    ph = Inches(2.0)
    gap = Inches(0.15)
    for i, (title, body, col) in enumerate(panels):
        x = px + (pw + gap) * i
        add_panel(s, x, py, pw, ph, border=col, border_w=Pt(1.5))
        add_text(s, x + Inches(0.3), py + Inches(0.25),
                 pw - Inches(0.6), Inches(0.5),
                 title, size=22, bold=True, color=col, font=FONT_TITLE)
        add_text(s, x + Inches(0.3), py + Inches(0.95),
                 pw - Inches(0.6), ph - Inches(1.1),
                 body, size=14, color=PARCHMENT)
    add_footer(s, 2, TOTAL)

    # ============================================================ 3. THEME FIT
    s = prs.slides.add_slide(blank)
    add_bg(s, "ink")
    add_kicker(s, Inches(0.7), Inches(0.55),
               "THEME  ·  ANCIENT TECHNOLOGY")
    add_text(s, Inches(0.7), Inches(0.95), Inches(11.5), Inches(1.0),
             "The Past Wrote The Future. Then Buried It.",
             size=34, bold=True, color=PARCHMENT_HI, font=FONT_TITLE)
    add_double_rule(s, Inches(0.75), Inches(1.85), Inches(2.0))

    add_text(s, Inches(0.7), Inches(2.10), Inches(12), Inches(1.4),
             "Babel's Glyph isn't *about* ancient technology \u2014 "
             "it *is* ancient technology.\n"
             "Every system, enemy, and zone is a forgotten machine "
             "clawing its way back to life.",
             size=18, color=PARCHMENT, italic=True, spacing=1.35)

    # Three pillars of the theme: PREMISE / SYSTEMS / WORLD
    px = Inches(0.7)
    py = Inches(3.85)
    pw = Inches(4.0)
    ph = Inches(2.95)
    gap = Inches(0.15)

    # 1. PREMISE
    add_panel(s, px, py, pw, ph, border=GLYPH_GOLD, border_w=Pt(1.5))
    add_text(s, px + Inches(0.3), py + Inches(0.20),
             pw - Inches(0.6), Inches(0.45),
             "PREMISE", size=18, bold=True, color=GLYPH_GOLD,
             font=FONT_TITLE)
    add_text(s, px + Inches(0.3), py + Inches(0.65),
             pw - Inches(0.6), Inches(0.5),
             "Glyphs as code", size=14, bold=True, color=PARCHMENT_HI,
             italic=True)
    add_paragraphs(s, px + Inches(0.3), py + Inches(1.20),
                   pw - Inches(0.6), ph - Inches(1.4),
                   [
                       "The Library of Babel held every machine",
                       "Glyphs are its programming language",
                       "You are a thief recompiling the past",
                   ], size=12, color=PARCHMENT, line_spacing=1.35,
                   bullet_glyph="\u25c6")

    # 2. SYSTEMS
    x2 = px + (pw + gap)
    add_panel(s, x2, py, pw, ph, border=EMBER, border_w=Pt(1.5))
    add_text(s, x2 + Inches(0.3), py + Inches(0.20),
             pw - Inches(0.6), Inches(0.45),
             "SYSTEMS", size=18, bold=True, color=EMBER,
             font=FONT_TITLE)
    add_text(s, x2 + Inches(0.3), py + Inches(0.65),
             pw - Inches(0.6), Inches(0.5),
             "Every mechanic is a relic", size=14, bold=True,
             color=PARCHMENT_HI, italic=True)
    add_paragraphs(s, x2 + Inches(0.3), py + Inches(1.20),
                   pw - Inches(0.6), ph - Inches(1.4),
                   [
                       "Glyph-bombs \u2014 arcane charges",
                       "Steam vents \u2014 Vitruvian pneumatics",
                       "Automatons \u2014 clockwork sentinels",
                       "Crumble tiles \u2014 decaying masonry",
                   ], size=12, color=PARCHMENT, line_spacing=1.35,
                   bullet_glyph="\u25c6")

    # 3. WORLD
    x3 = px + 2 * (pw + gap)
    add_panel(s, x3, py, pw, ph, border=LAPIS_LIGHT, border_w=Pt(1.5))
    add_text(s, x3 + Inches(0.3), py + Inches(0.20),
             pw - Inches(0.6), Inches(0.45),
             "WORLD", size=18, bold=True, color=LAPIS_LIGHT,
             font=FONT_TITLE)
    add_text(s, x3 + Inches(0.3), py + Inches(0.65),
             pw - Inches(0.6), Inches(0.5),
             "Three eras of lost engineering", size=14, bold=True,
             color=PARCHMENT_HI, italic=True)
    add_paragraphs(s, x3 + Inches(0.3), py + Inches(1.20),
                   pw - Inches(0.6), ph - Inches(1.4),
                   [
                       "Sandstone \u2014 pre-mechanical glyph-craft",
                       "Forge \u2014 Renaissance proto-machines",
                       "Sky workshop \u2014 impossible engines",
                   ], size=12, color=PARCHMENT, line_spacing=1.35,
                   bullet_glyph="\u25c6")

    # Bottom tagline ribbon
    add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.35),
             "\u201cYou are not running away from the past. "
             "You are running through it \u2014 and rewriting it as you go.\u201d",
             size=12, color=GLYPH_GOLD, italic=True,
             align=PP_ALIGN.CENTER, font=FONT_TITLE)
    add_footer(s, 3, TOTAL)
    s = prs.slides.add_slide(blank)
    add_bg(s, "parchment")
    add_kicker(s, Inches(0.7), Inches(0.55), "CORE PILLARS",
               color=COPPER)
    add_text(s, Inches(0.7), Inches(0.95), Inches(11), Inches(1.0),
             "Three Pillars Hold The Game Up",
             size=38, bold=True, color=INK, font=FONT_TITLE)
    add_double_rule(s, Inches(0.75), Inches(1.85), Inches(2.0),
                    color=COPPER)

    pillars = [
        ("01", "MOVEMENT", COPPER,
         "Tight, expressive, modern.",
         ["Double-jump with cut-off",
          "Dash with i-frames",
          "Wall-slide & wall-jump",
          "Slide under low ceilings",
          "Coyote time + jump buffer"]),
        ("02", "WORLD", LAPIS,
         "Three zones. One forbidden library.",
         ["Sandstone outskirts",
          "Da Vinci's burning forge",
          "Workshop above the clouds",
          "Parallax skies, runtime tinting",
          "30-tile authored chunks"]),
        ("03", "THREAT", EMBER_DIM,
         "The past is hunting you.",
         ["Spike traps",
          "Crumble tiles",
          "Steam-vent launchers",
          "Patrolling automatons",
          "Glyph-bomb counterplay"]),
    ]
    px = Inches(0.7)
    py = Inches(2.15)
    pw = Inches(4.0)
    ph = Inches(4.65)
    gap = Inches(0.15)
    for i, (num, title, col, sub, items) in enumerate(pillars):
        x = px + (pw + gap) * i
        add_panel(s, x, py, pw, ph, border=col, border_w=Pt(1.5))
        add_text(s, x + Inches(0.35), py + Inches(0.25),
                 Inches(1.5), Inches(0.7),
                 num, size=44, bold=True, color=col, font=FONT_TITLE)
        add_text(s, x + Inches(0.35), py + Inches(1.10),
                 pw - Inches(0.7), Inches(0.5),
                 title, size=20, bold=True, color=INK, font=FONT_TITLE)
        add_text(s, x + Inches(0.35), py + Inches(1.55),
                 pw - Inches(0.7), Inches(0.5),
                 sub, size=13, color=INK_SOFT, italic=True)
        add_paragraphs(s, x + Inches(0.35), py + Inches(2.20),
                       pw - Inches(0.7), Inches(2.4),
                       items, size=13, color=INK_SOFT, line_spacing=1.35,
                       bullet_glyph="◆")
    add_footer(s, 4, TOTAL, color=INK_FAINT)

    # ============================================================ 5. MOVEMENT
    s = prs.slides.add_slide(blank)
    add_bg(s, "parchment")
    add_kicker(s, Inches(0.7), Inches(0.55), "FEATURE  ·  MOVEMENT",
               color=COPPER)
    add_text(s, Inches(0.7), Inches(0.95), Inches(11), Inches(1.0),
             "Movement Is The Whole Game",
             size=38, bold=True, color=INK, font=FONT_TITLE)
    add_double_rule(s, Inches(0.75), Inches(1.85), Inches(2.0),
                    color=COPPER)

    add_text(s, Inches(0.7), Inches(2.10), Inches(12), Inches(0.6),
             "Every input is one frame off muscle-memory. "
             "Nothing is decorative.",
             size=18, color=INK_SOFT, italic=True)

    add_panel(s, Inches(0.7), Inches(2.95), Inches(5.9), Inches(4.0),
              border=COPPER, border_w=Pt(1.5))
    add_text(s, Inches(1.0), Inches(3.10), Inches(5.5), Inches(0.5),
             "VERBS", size=20, bold=True, color=COPPER, font=FONT_TITLE)
    add_paragraphs(s, Inches(1.0), Inches(3.65),
                   Inches(5.5), Inches(3.2),
                   [
                       ("Run · auto-scroll keeps the pressure on",
                        INK, True),
                       "Single-jump with variable height",
                       "Double-jump with cut-off",
                       ("Dash · short i-frame burst, kills enemies",
                        INK, True),
                       "Slide · drops the hitbox + boosts momentum",
                       ("Wall-slide & wall-jump",
                        INK, True),
                       ("Glyph-bomb · throwable arcane charge",
                        EMBER_DIM, True),
                   ], size=14, color=INK_SOFT)

    add_panel(s, Inches(6.75), Inches(2.95), Inches(5.9), Inches(4.0),
              border=LAPIS, border_w=Pt(1.5))
    add_text(s, Inches(7.05), Inches(3.10), Inches(5.5), Inches(0.5),
             "GAME FEEL", size=20, bold=True, color=LAPIS,
             font=FONT_TITLE)
    add_paragraphs(s, Inches(7.05), Inches(3.65),
                   Inches(5.5), Inches(3.2),
                   [
                       "Coyote time — forgive late jumps",
                       "Jump buffering — forgive early jumps",
                       "Air control with momentum carry",
                       "Hit-stop and screen-shake on impact",
                       "Stomp refunds an air-jump",
                       "Particle bursts on every action",
                       ("Locked 60 FPS, fixed-step physics",
                        INK, True),
                   ], size=14, color=INK_SOFT)

    add_footer(s, 5, TOTAL, color=INK_FAINT)

    # ============================================================ 6. ZONES
    s = prs.slides.add_slide(blank)
    add_bg(s, "ink")
    add_kicker(s, Inches(0.7), Inches(0.55), "FEATURE  ·  WORLD")
    add_text(s, Inches(0.7), Inches(0.95), Inches(11), Inches(1.0),
             "Three Zones. Every One Wants You Dead.",
             size=34, bold=True, color=PARCHMENT_HI, font=FONT_TITLE)
    add_double_rule(s, Inches(0.75), Inches(1.85), Inches(2.0))

    # Each zone panel uses ITS in-game palette.
    zones = [
        ("ZONE  I",   "SANDSTONE OUTSKIRTS",
         SAND_LIGHT, SAND_DARK, INK,
         "0 m — 3 km",
         "Crumbling ruins. Spike pits. Where overconfident new players die."),
        ("ZONE  II",  "DA VINCI'S FORGE",
         EMBER, SKY_FORGE_TOP, PARCHMENT_HI,
         "3 km — 8 km",
         "Iron beams, steam-vent launchers, patrolling automatons. "
         "Vertical play opens up. The tempo doubles."),
        ("ZONE  III", "WORKSHOP IN THE SKY",
         SKY_WORK_TOP, LAPIS, INK,
         "8 km +",
         "Wind, gaps, razor-thin platforms. Wall-runs are mandatory. "
         "Gravity gets opinionated."),
    ]
    px = Inches(0.7)
    py = Inches(2.20)
    pw = Inches(4.0)
    ph = Inches(4.55)
    gap = Inches(0.15)
    for i, (tag, name, fill, border, ink_col, dist, body) in enumerate(zones):
        x = px + (pw + gap) * i
        add_panel(s, x, py, pw, ph, fill=fill, border=border,
                  border_w=Pt(1.8))
        add_text(s, x + Inches(0.35), py + Inches(0.25),
                 pw - Inches(0.7), Inches(0.4),
                 tag, size=12, bold=True, color=border, font=FONT_TITLE)
        add_text(s, x + Inches(0.35), py + Inches(0.70),
                 pw - Inches(0.7), Inches(1.4),
                 name, size=22, bold=True, color=ink_col,
                 font=FONT_TITLE)
        add_text(s, x + Inches(0.35), py + Inches(1.85),
                 pw - Inches(0.7), Inches(0.4),
                 dist, size=13, color=border, bold=True,
                 font=FONT_MONO)
        add_text(s, x + Inches(0.35), py + Inches(2.35),
                 pw - Inches(0.7), ph - Inches(2.6),
                 body, size=14, color=ink_col)
    add_footer(s, 6, TOTAL)

    # ============================================================ 7. SYSTEMIC
    s = prs.slides.add_slide(blank)
    add_bg(s, "ink")
    add_kicker(s, Inches(0.7), Inches(0.55),
               "FEATURE  ·  SYSTEMIC DESIGN")
    add_text(s, Inches(0.7), Inches(0.95), Inches(11), Inches(1.0),
             "Nothing Is Hand-Placed.",
             size=42, bold=True, color=PARCHMENT_HI, font=FONT_TITLE)
    add_double_rule(s, Inches(0.75), Inches(1.95), Inches(2.0))

    add_text(s, Inches(0.7), Inches(2.20), Inches(12), Inches(1.0),
             "The world is stitched together at runtime from a library "
             "of small,\n"
             "hand-authored 30-tile chunks — picked, biased, and threaded "
             "so every run is fresh.",
             size=18, color=PARCHMENT, spacing=1.35, italic=True)

    # Chunk pipeline — five chained scrolls of parchment with glyph-arrows
    flow_y = Inches(3.80)
    chunk_w = Inches(1.65)
    chunk_h = Inches(1.45)
    chunk_gap = Inches(0.55)
    chunk_x = Inches(0.7)
    labels = ["CHUNK\n#A12", "CHUNK\n#B07", "CHUNK\n#C24",
              "CHUNK\n#A18", "CHUNK\n#D03"]
    border_colors = [GLYPH_GOLD, LAPIS_LIGHT, EMBER, GLYPH_GOLD,
                     LAPIS_LIGHT]
    for i, (lbl, col) in enumerate(zip(labels, border_colors)):
        x = chunk_x + i * (chunk_w + chunk_gap)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x, flow_y, chunk_w, chunk_h)
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(PARCHMENT)
        box.line.color.rgb = rgb(col)
        box.line.width = Pt(1.6)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = lbl
        r.font.name = FONT_MONO
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = rgb(INK)

        if i < len(labels) - 1:
            ax = x + chunk_w
            ay = flow_y + chunk_h / 2
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       ax + Inches(0.08),
                                       ay - Inches(0.13),
                                       chunk_gap - Inches(0.16),
                                       Inches(0.26))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb(GLYPH_GOLD)
            arrow.line.fill.background()

    add_text(s, Inches(0.7), flow_y + chunk_h + Inches(0.30),
             Inches(12), Inches(0.4),
             "Picker biases by zone & distance  ·  entrances/exits "
             "auto-align  ·  scroll speed ramps with km traveled",
             size=12, color=INK_FAINT, align=PP_ALIGN.CENTER,
             font=FONT_MONO)

    # Side stat
    sx = Inches(10.4)
    sy = Inches(2.10)
    add_panel(s, sx, sy, Inches(2.5), Inches(1.5), border=GLYPH_GOLD,
              border_w=Pt(1.5))
    add_text(s, sx + Inches(0.2), sy + Inches(0.18),
             Inches(2.2), Inches(0.4),
             "PICKER POOL", size=11, bold=True, color=GLYPH_GOLD,
             font=FONT_TITLE)
    add_text(s, sx + Inches(0.2), sy + Inches(0.50),
             Inches(2.2), Inches(0.95),
             "30+", size=46, bold=True, color=PARCHMENT_HI,
             font=FONT_TITLE)
    add_footer(s, 7, TOTAL)

    # ============================================================ 8. TECH
    s = prs.slides.add_slide(blank)
    add_bg(s, "ink")
    add_kicker(s, Inches(0.7), Inches(0.55), "TECH")
    add_text(s, Inches(0.7), Inches(0.95), Inches(11), Inches(1.0),
             "Built From Scratch In Python.",
             size=40, bold=True, color=PARCHMENT_HI, font=FONT_TITLE)
    add_double_rule(s, Inches(0.75), Inches(1.95), Inches(2.0))

    add_text(s, Inches(0.7), Inches(2.20), Inches(12), Inches(0.6),
             "No engine. No store-bought assets. The constraint is the design.",
             size=18, color=PARCHMENT, italic=True)

    add_panel(s, Inches(0.7), Inches(3.05), Inches(5.9), Inches(3.85),
              border=GLYPH_GOLD, border_w=Pt(1.5))
    add_text(s, Inches(1.0), Inches(3.20), Inches(5.5), Inches(0.5),
             "STACK", size=20, bold=True, color=GLYPH_GOLD,
             font=FONT_TITLE)
    add_paragraphs(s, Inches(1.0), Inches(3.75),
                   Inches(5.5), Inches(3.0),
                   [
                       ("Python 3.12  ·  pygame-ce", PARCHMENT_HI, True),
                       "Custom 2-D engine: render, physics, AI, audio",
                       "Fixed-step deterministic simulation",
                       "All sprites rasterised in code at runtime",
                       "Procedural music & SFX with NumPy",
                       ("PyInstaller → single .exe (one build.bat)",
                        GLYPH_GOLD, True),
                   ], size=14, color=PARCHMENT, line_spacing=1.35)

    add_panel(s, Inches(6.75), Inches(3.05), Inches(5.9), Inches(3.85),
              border=LAPIS_LIGHT, border_w=Pt(1.5))
    add_text(s, Inches(7.05), Inches(3.20), Inches(5.5), Inches(0.5),
             "ARCHITECTURE", size=20, bold=True, color=LAPIS_LIGHT,
             font=FONT_TITLE)
    add_paragraphs(s, Inches(7.05), Inches(3.75),
                   Inches(5.5), Inches(3.0),
                   [
                       "world.py    ·  scrolling tilemap & camera",
                       "chunks.py   ·  authored rooms + picker",
                       "player.py   ·  state machine, all 7 verbs",
                       "entities.py ·  hazards, foes, glyph-bombs",
                       "particles.py ·  fire / steam / dust",
                       "hud.py · render.py · input.py",
                   ], size=13, color=PARCHMENT, line_spacing=1.4,
                   font=FONT_MONO, bullet_glyph="›")

    add_footer(s, 8, TOTAL)

    # ============================================================ 9. NUMBERS
    s = prs.slides.add_slide(blank)
    add_bg(s, "ink")
    add_kicker(s, Inches(0.7), Inches(0.55), "BY THE NUMBERS")
    add_text(s, Inches(0.7), Inches(0.95), Inches(11), Inches(1.0),
             "The Game In Numbers",
             size=40, bold=True, color=PARCHMENT_HI, font=FONT_TITLE)
    add_double_rule(s, Inches(0.75), Inches(1.95), Inches(2.0))

    stats = [
        ("60",     "FPS",            "locked, fixed-step",  GLYPH_GOLD),
        ("2,900+", "LINES OF PYTHON","across 9 modules",    LAPIS_LIGHT),
        ("3",      "ZONES",          "30+ chunk variants",  EMBER),
        ("7",      "PLAYER VERBS",   "all combo-able",      GLYPH_GOLD),
        ("1",      "EXECUTABLE",     "via PyInstaller",     LAPIS_LIGHT),
        ("30 s",   "PROMO TRAILER",  "Azure neural VO",     EMBER),
    ]
    cols = 3
    pw = Inches(4.0)
    ph = Inches(2.10)
    gap_x = Inches(0.15)
    gap_y = Inches(0.20)
    px0 = Inches(0.7)
    py0 = Inches(2.25)
    for i, (n, lbl, sub, col) in enumerate(stats):
        r, c = divmod(i, cols)
        x = px0 + c * (pw + gap_x)
        y = py0 + r * (ph + gap_y)
        add_panel(s, x, y, pw, ph, border=col, border_w=Pt(1.5))
        add_text(s, x + Inches(0.3), y + Inches(0.18),
                 pw - Inches(0.6), Inches(1.2),
                 n, size=54, bold=True, color=col, font=FONT_TITLE)
        add_text(s, x + Inches(0.3), y + Inches(1.30),
                 pw - Inches(0.6), Inches(0.4),
                 lbl, size=14, bold=True, color=PARCHMENT_HI,
                 font=FONT_TITLE)
        add_text(s, x + Inches(0.3), y + Inches(1.65),
                 pw - Inches(0.6), Inches(0.4),
                 sub, size=11, color=INK_FAINT, italic=True)
    add_footer(s, 9, TOTAL)

    # ============================================================ 10. ROADMAP
    s = prs.slides.add_slide(blank)
    add_bg(s, "sky")
    add_kicker(s, Inches(0.7), Inches(0.55), "ROADMAP", color=LAPIS)
    add_text(s, Inches(0.7), Inches(0.95), Inches(11), Inches(1.0),
             "Where We Go Next",
             size=42, bold=True, color=INK, font=FONT_TITLE)
    add_double_rule(s, Inches(0.75), Inches(1.95), Inches(2.0),
                    color=LAPIS)

    horizons = [
        ("NEAR",  "next 1 — 2 months", COPPER, [
            "Daily seed challenge & leaderboards",
            "Two new chunk packs per zone",
            "Settings menu — rebind, audio sliders",
            "Save profile + meta-progression unlocks",
        ]),
        ("MID",   "next quarter",      LAPIS, [
            "Fourth zone — submerged crypt",
            "Boss encounters at zone transitions",
            "New traversal verb: grapple-glyph",
            "Steam Deck-friendly build & controller polish",
        ]),
        ("FAR",   "stretch goals",     EMBER_DIM, [
            "Online ghost races (async multiplayer)",
            "Mobile (Android/iOS) via pygame-ce + Buildozer",
            "Mod tools — author & share chunk packs",
            "Procedural soundtrack reactive to distance",
        ]),
    ]
    px = Inches(0.7)
    py = Inches(2.20)
    pw = Inches(4.0)
    ph = Inches(4.6)
    gap = Inches(0.15)
    for i, (tag, when, col, items) in enumerate(horizons):
        x = px + (pw + gap) * i
        add_panel(s, x, py, pw, ph, border=col, border_w=Pt(1.5))
        add_text(s, x + Inches(0.35), py + Inches(0.30),
                 pw - Inches(0.7), Inches(0.6),
                 tag, size=24, bold=True, color=col, font=FONT_TITLE)
        add_text(s, x + Inches(0.35), py + Inches(0.95),
                 pw - Inches(0.7), Inches(0.4),
                 when.upper(), size=11, bold=True, color=INK_FAINT,
                 font=FONT_MONO)
        add_paragraphs(s, x + Inches(0.35), py + Inches(1.55),
                       pw - Inches(0.7), ph - Inches(1.7),
                       items, size=14, color=INK_SOFT,
                       line_spacing=1.4, bullet_glyph="◆")
    add_footer(s, 10, TOTAL, color=INK_FAINT)

    # ============================================================ 11. CLOSE
    s = prs.slides.add_slide(blank)
    add_bg(s, "title")

    add_text(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.5),
             "THANK YOU", size=18, bold=True, color=COPPER,
             font=FONT_TITLE)

    add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.4),
             "Run far. Die loud.", size=78, bold=True, color=INK,
             font=FONT_TITLE)
    add_text(s, Inches(0.9), Inches(3.20), Inches(11.5), Inches(1.4),
             "Rewrite the past.", size=78, bold=True, color=INK,
             font=FONT_TITLE, italic=True)

    add_double_rule(s, Inches(0.95), Inches(4.85), Inches(2.4),
                    color=COPPER)

    add_text(s, Inches(0.9), Inches(5.10), Inches(11.5), Inches(0.6),
             "Babel's Glyph  ·  built in Python  ·  questions?",
             size=22, color=INK_SOFT, font=FONT_TITLE, italic=True)

    add_text(s, Inches(0.9), Inches(6.40), Inches(11.5), Inches(0.4),
             "github.com/<your-handle>/BabelsGlyph",
             size=14, color=INK_FAINT, font=FONT_MONO)

    try:
        prs.save(OUT_PATH)
        return OUT_PATH
    except PermissionError:
        # File is open in PowerPoint — fall back to a sibling filename so
        # the user can compare before replacing.
        alt = OUT_PATH.with_name(OUT_PATH.stem + "_new" + OUT_PATH.suffix)
        prs.save(alt)
        print(f"[warn] {OUT_PATH.name} is open / locked; wrote {alt.name} "
              f"instead. Close PowerPoint and re-run, or replace manually.")
        return alt


if __name__ == "__main__":
    out = build()
    size_kb = out.stat().st_size / 1024
    print(f"Wrote {out}  ({size_kb:.0f} KB)")
